from dataclasses import dataclass
from pathlib import Path
from typing import Any

from multimodal_rag.utils.logging_utils import logging

logger = logging.getLogger(__name__)

_SUPPORTED_OFFICE_EXTS = frozenset({".docx", ".pptx", ".odt", ".odp"})


@dataclass
class OfficeProcessor:
    """Extract text and images from Office documents.

    Supports ``.docx`` (Word), ``.pptx`` (PowerPoint), ``.odt`` (Writer),
    and ``.odp`` (Impress) formats.

    * Word documents are chunked by paragraph groups.
    * Presentations are chunked per slide.
    * ODF documents are text-only (no image extraction).
    """

    chunk_size: int = 8192
    chunk_overlap: int = 512
    text_splitter: Any | None = None

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def process(self, office_path: str) -> list[dict[str, Any]]:
        ext = Path(office_path).suffix.lower()
        source = str(office_path)

        if ext == ".docx":
            return self._process_docx(office_path, source)
        elif ext == ".pptx":
            return self._process_pptx(office_path, source)
        elif ext == ".odt":
            return self._process_odt(office_path, source)
        elif ext == ".odp":
            return self._process_odp(office_path, source)
        else:
            raise ValueError(f"Unsupported Office format: {ext}")

    # ------------------------------------------------------------------
    # DOCX
    # ------------------------------------------------------------------

    def _process_docx(self, path: str, source: str) -> list[dict[str, Any]]:
        try:
            from docx import Document as DocxDocument
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
        except ImportError:
            raise ImportError("python-docx is required for .docx files. Install with: pip install python-docx")

        try:
            doc = DocxDocument(path)
        except Exception as e:
            logger.warning("Failed to open .docx %s: %s", path, e)
            return []

        # Extract inline images keyed by relationship ID
        images: dict[str, str] = {}
        for rel in doc.part.rels.values():
            if rel.reltype == RT.IMAGE:
                try:
                    img_data = rel.target_part.blob
                    mime = rel.target_part.content_type
                    import base64

                    b64 = base64.b64encode(img_data).decode("utf-8")
                    images[rel.target_part.partname] = f"data:{mime};base64,{b64}"
                except Exception:
                    logger.debug("Suppressed exception", exc_info=True)

        # Extract paragraphs with their inline images
        paragraphs: list[dict[str, Any]] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            para_images: list[str] = []
            for run in para.runs:
                for inline in run._element.findall(
                    ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                ):
                    blip = inline.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip")
                    for b in blip:
                        embed = b.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                        if embed:
                            partname = doc.part.rels[embed].target_part.partname
                            if partname in images:
                                para_images.append(images[partname])
            paragraphs.append({"text": text, "images": para_images})

        # Also extract tables
        for table in doc.tables:
            rows_text = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows_text.append(" | ".join(cells))
            if rows_text:
                paragraphs.append({"text": "\n".join(rows_text), "images": []})

        if not paragraphs:
            return []

        return self._chunk_paragraphs(paragraphs, source, "docx")

    # ------------------------------------------------------------------
    # PPTX
    # ------------------------------------------------------------------

    def _process_pptx(self, path: str, source: str) -> list[dict[str, Any]]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError("python-pptx is required for .pptx files. Install with: pip install python-pptx")

        try:
            prs = Presentation(path)
        except Exception as e:
            logger.warning("Failed to open .pptx %s: %s", path, e)
            return []

        slides: list[dict[str, Any]] = []
        for slide_idx, slide in enumerate(prs.slides):
            slide_texts: list[str] = []
            slide_images: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img_data = shape.image.blob
                        mime = shape.image.content_type
                        import base64

                        b64 = base64.b64encode(img_data).decode("utf-8")
                        slide_images.append(f"data:{mime};base64,{b64}")
                    except Exception:
                        logger.debug("Suppressed exception", exc_info=True)

            text = "\n".join(slide_texts) if slide_texts else ""
            if text or slide_images:
                entry: dict[str, Any] = {
                    "text": text,
                    "source": source,
                    "slide": slide_idx + 1,
                }
                if slide_images:
                    entry["image"] = slide_images
                slides.append(entry)

        return slides if slides else []

    # ------------------------------------------------------------------
    # ODT
    # ------------------------------------------------------------------

    @staticmethod
    def _process_odt(path: str, source: str) -> list[dict[str, Any]]:
        try:
            from odf import table, text
            from odf.opendocument import load
        except ImportError:
            raise ImportError("odfpy is required for .odt files. Install with: pip install odfpy")

        try:
            doc = load(path)
        except Exception as e:
            logger.warning("Failed to open .odt %s: %s", path, e)
            return []

        texts: list[str] = []

        paragraphs = doc.getElementsByType(text.P)
        for p in paragraphs:
            t = _odf_text_content(p)
            if t:
                texts.append(t)

        tables = doc.getElementsByType(table.Table)
        for tbl in tables:
            rows_text: list[str] = []
            for row in tbl.getElementsByType(table.TableRow):
                cells = []
                for cell in row.getElementsByType(table.TableCell):
                    cell_texts = []
                    for p in cell.getElementsByType(text.P):
                        ct = _odf_text_content(p)
                        if ct:
                            cell_texts.append(ct)
                    cells.append(" ".join(cell_texts))
                if cells:
                    rows_text.append(" | ".join(cells))
            if rows_text:
                texts.append("\n".join(rows_text))

        if not texts:
            return []

        return [{"text": "\n\n".join(texts), "source": source}]

    # ------------------------------------------------------------------
    # ODP
    # ------------------------------------------------------------------

    @staticmethod
    def _process_odp(path: str, source: str) -> list[dict[str, Any]]:
        try:
            from odf import draw
            from odf import text as odf_text
            from odf.opendocument import load
        except ImportError:
            raise ImportError("odfpy is required for .odp files. Install with: pip install odfpy")

        try:
            doc = load(path)
        except Exception as e:
            logger.warning("Failed to open .odp %s: %s", path, e)
            return []

        slides: list[dict[str, Any]] = []
        for i, page in enumerate(doc.getElementsByType(draw.Page)):
            slide_texts: list[str] = []
            for p in page.getElementsByType(odf_text.P):
                t = _odf_text_content(p)
                if t:
                    slide_texts.append(t)
            if slide_texts:
                slides.append({"text": "\n".join(slide_texts), "source": source, "slide": i + 1})

        return slides

    # ------------------------------------------------------------------
    # Budget helpers
    # ------------------------------------------------------------------

    def _exceeds_budget(self, text: str) -> bool:
        if self.text_splitter is not None:
            return self.text_splitter.count_tokens(text) > self.chunk_size
        return len(text) > self.chunk_size

    def _exceeds_combined_budget(self, a: str, b: str) -> bool:
        if self.text_splitter is not None:
            combined = a + "\n\n" + b if a and b else a or b
            return self.text_splitter.count_tokens(combined) > self.chunk_size
        return len(a) + len(b) + 1 > self.chunk_size

    # ------------------------------------------------------------------
    # Chunking
    # ------------------------------------------------------------------

    def _chunk_paragraphs(self, paragraphs: list[dict[str, Any]], source: str, fmt: str) -> list[dict[str, Any]]:
        chunks: list[dict[str, Any]] = []
        buffer_text = ""
        buffer_images: list[tuple[str, int]] = []

        def flush() -> None:
            nonlocal buffer_text, buffer_images
            if buffer_text.strip():
                chunk: dict[str, Any] = {"text": buffer_text.strip(), "source": source}
                imgs = list(dict.fromkeys(img for img, _ in buffer_images))
                if imgs:
                    chunk["image"] = imgs
                chunks.append(chunk)
            buffer_text = ""
            buffer_images = []

        for para in paragraphs:
            p_text = para["text"]
            p_imgs = [(img, id(img)) for img in para.get("images", [])]

            if not buffer_text:
                buffer_text = p_text
                buffer_images = p_imgs
            elif self._exceeds_combined_budget(buffer_text, p_text):
                flush()
                buffer_text = p_text
                buffer_images = p_imgs
            else:
                buffer_text += "\n\n" + p_text
                buffer_images.extend(p_imgs)

        flush()
        return chunks


def _odf_text_content(element) -> str:
    """Extract text content from an ODF element, including inside ``<text:span>`` etc."""
    texts: list[str] = []
    for node in element.childNodes:
        if node.nodeType == node.TEXT_NODE:
            texts.append(node.data)
        elif hasattr(node, "getElementsByType"):
            from odf import text

            for child in node.getElementsByType(text.P):
                texts.append(_odf_text_content(child))
            for child in node.getElementsByType(text.Span):
                texts.append(_odf_text_content(child))
            for child in node.getElementsByType(text.S):
                texts.append(" ")
            for child in node.getElementsByType(text.Tab):
                texts.append("\t")
            for child in node.getElementsByType(text.LineBreak):
                texts.append("\n")
    return "".join(texts).strip()
